#!/usr/bin/env python3
"""
IMA知识库 → Obsidian 批量下载+转Markdown
从已有 .meta.json + .url.txt 出发，下载原文并提取文字生成 .md 笔记
"""
import json, os, subprocess, sys, glob, re, time
from pathlib import Path
from datetime import datetime, timezone

HOME = os.environ['HOME']
SKILL_DIR = os.path.join(HOME, '.hermes/skills/ima')
IMA_API = os.path.join(SKILL_DIR, 'ima_api.cjs')
KB_SYNC_DIR = os.path.join(HOME, 'knowledge/ima-sync')
DOWNLOAD_DIR = os.path.join(KB_SYNC_DIR, 'downloads')
OUTPUT_DIRS = {
    'telecom': os.path.join(HOME, 'knowledge/telecom/ima-articles'),
    'general': os.path.join(HOME, 'knowledge/ima-sync/notes'),
}
STATE_FILE = os.path.join(KB_SYNC_DIR, '.download_state.json')
MEDIA_TYPE_NAMES = {1:'PDF',2:'图片',3:'DOC',4:'PPT',5:'视频',6:'网页链接',7:'Excel',8:'TXT',11:'笔记',12:'思维导图'}

def load_state():
    if os.path.exists(STATE_FILE):
        with open(STATE_FILE) as f: return json.load(f)
    return {"downloaded_ids": [], "converted_ids": []}

def save_state(state):
    state['last_run'] = datetime.now(timezone.utc).isoformat()
    with open(STATE_FILE, 'w') as f: json.dump(state, f, ensure_ascii=False, indent=2)

def call_ima(api_path, body):
    cmd = ['node', IMA_API, api_path, json.dumps(body)]
    try:
        r = subprocess.run(cmd, cwd=SKILL_DIR, capture_output=True, text=True, timeout=30)
        return json.loads(r.stdout) if r.returncode == 0 else None
    except: return None

def download_file(url, headers, dest):
    os.makedirs(os.path.dirname(dest), exist_ok=True)
    if os.path.exists(dest) and os.path.getsize(dest) > 0: return True
    curl = ['curl', '-sL', '-o', dest, '--connect-timeout', '15', '--max-time', '60']
    for k,v in headers.items(): curl.extend(['-H', f"{k}: {v}"])
    curl.append(url)
    try: subprocess.run(curl, check=True, timeout=120); return os.path.getsize(dest) > 0
    except: return False

def extract_text(path, media_type):
    ext = os.path.splitext(path)[1].lower()
    if media_type == 1 or ext == '.pdf':
        txt = path + '.txt'
        subprocess.run(['pdftotext', '-layout', path, txt], capture_output=True, timeout=60)
        if os.path.exists(txt):
            with open(txt) as f: return f.read()
    elif media_type in (3, 4) or ext in ('.docx', '.pptx'):
        try:
            if media_type == 4 or ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(path)
                return '\n\n'.join(s.text for s in prs.slides for sh in s.shapes if hasattr(sh,'text') and sh.text.strip())
            else:
                import docx
                d = docx.Document(path)
                return '\n'.join(p.text for p in d.paragraphs)
        except: pass
    elif media_type == 7 or ext in ('.xls','.xlsx'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for sn in wb.sheetnames:
                ws = wb[sn]
                cells = [' | '.join(str(c) for c in r if c) for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)]
                if cells: rows.append(f"=== {sn} ===\n" + '\n'.join(cells[:200]))
            return '\n\n'.join(rows)
        except: pass
    elif media_type == 8 or ext == '.txt':
        for enc in ['utf-8','gbk','gb2312','big5','latin-1']:
            try:
                with open(path, encoding=enc) as f: return f.read()
            except: continue
    return ""

def categorize(text, title):
    li_kw = ['合法监听','LI','Lawful','X1','X2','X3','监听','LIID']
    ims_kw = ['IMS','VoLTE','VoNR','SIP','SBC','CSCF','MMTel','P-CSCF']
    g_kw = ['5G','5GC','NR','NGAP','AMF','SMF','UPF','PCF','NRF','NSSF','网络切片','S-NSSAI']
    for kw in li_kw:
        if kw.lower() in title.lower() or kw.lower() in text[:500].lower(): return 'telecom'
    for kw in ims_kw:
        if kw.lower() in title.lower(): return 'telecom'
    for kw in g_kw:
        if kw.lower() in title.lower(): return 'telecom'
    return 'general'

def generate_md(title, text, meta):
    cat = categorize(text, title)
    out_dir = OUTPUT_DIRS.get(cat, OUTPUT_DIRS['general'])
    os.makedirs(out_dir, exist_ok=True)
    safe = re.sub(r'[<>:"/\\|?*]', '', title).strip()[:200]
    for ext in ['.pdf','.doc','.docx','.ppt','.pptx','.xlsx','.txt']:
        if safe.endswith(ext): safe = safe[:-len(ext)]
    safe = safe.strip()
    md_path = os.path.join(out_dir, f"{safe}.md")
    body = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text).strip()[:15000]
    mt = MEDIA_TYPE_NAMES.get(meta.get('media_type',0),'未知')
    content = f"""---
title: "{safe}"
source: "IMA知识库 - {meta.get('kb_name','未知')}"
media_type: "{mt}"
media_id: "{meta.get('media_id','')}"
synced_at: "{meta.get('synced_at','')}"
category: "{cat}"
---
# {safe}
> 📥 IMA知识库「{meta.get('kb_name','未知')}」· {mt}
---
{body}
"""
    with open(md_path, 'w', encoding='utf-8') as f: f.write(content)
    return md_path

def main():
    state = load_state()
    downloaded = set(state.get('downloaded_ids',[]))
    converted = set(state.get('converted_ids',[]))
    metas = [m for m in glob.glob(f'{KB_SYNC_DIR}/**/*.meta.json', recursive=True) if '/downloads/' not in m]
    new_dl = new_cv = fail = 0
    for i,mp in enumerate(metas):
        with open(mp) as f: meta = json.load(f)
        mid = meta.get('media_id','')
        mt = meta.get('media_type',0)
        if mt in (2,5,6,12):
            continue
        if mid in downloaded and mid in converted:
            continue
        if mid not in downloaded:
            resp = call_ima('openapi/wiki/v1/get_media_info', {"knowledge_base_id":meta['kb_id'],"media_id":mid})
            if not resp or resp.get('code') != 0:
                fail += 1; continue
            ui = resp.get('data',{}).get('url_info',{})
            url = ui.get('url','')
            if not url:
                fail += 1; continue
            ext = os.path.splitext(meta.get('title',''))[1] or '.bin'
            dest = os.path.join(DOWNLOAD_DIR, meta.get('kb_name','unknown'), f"{mid}{ext}")
            if download_file(url, ui.get('headers',{}), dest):
                downloaded.add(mid); new_dl += 1
            else:
                fail += 1; continue
        else:
            # find file
            kb_dir = os.path.join(DOWNLOAD_DIR, meta.get('kb_name','unknown'))
            dest = ''
            if os.path.exists(kb_dir):
                for fn in os.listdir(kb_dir):
                    if mid in fn: dest = os.path.join(kb_dir, fn); break
        if mid not in converted and dest and os.path.getsize(dest) > 0:
            text = extract_text(dest, mt)
            if text.strip():
                generate_md(meta.get('title',mid), text, meta)
                converted.add(mid); new_cv += 1
        if (i+1) % 20 == 0:
            state['downloaded_ids'] = list(downloaded); state['converted_ids'] = list(converted)
            save_state(state)
        time.sleep(0.3)
    state['downloaded_ids'] = list(downloaded); state['converted_ids'] = list(converted)
    save_state(state)
    print(f"📊 下载:{new_dl} 转换:{new_cv} 失败:{fail} 累计:{len(downloaded)}/{len(converted)}")

if __name__ == '__main__':
    main()
